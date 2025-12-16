"""
新版PPT生成器 - 使用合并后的统一模板
完全匹配实际PPT模板的制作逻辑
"""
import os
import json
from pptx import Presentation

import config
import utils
from ai_image_generator import (
    generate_cover_image,
    generate_lecture_title_image,
    generate_intro_image,
    generate_knowledge_point_image,
    generate_learning_objectives_image
)
from slide_builder import SlideBuilder


def load_course_data():
    """加载课程数据"""
    if not os.path.exists(config.JSON_PATH):
        print(f"❌ 错误: 找不到数据文件 {config.JSON_PATH}")
        print("请先运行: python Smart_PPT_Factory/parser.py")
        return None
    
    with open(config.JSON_PATH, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    return data


def get_cover_info():
    """获取封面信息"""
    import glob
    
    # 尝试从PDF文件名解析
    pdfs = glob.glob(os.path.join(config.PDF_DIR, "*.pdf"))
    if pdfs:
        pdf_path = pdfs[0]
        cover_info = utils.parse_filename_to_json(pdf_path)
        return cover_info
    
    # 默认信息
    return {
        "subject": "语文",
        "grade": "高一",
        "season": "寒假",
        "teacher": "XXX老师",
        "subtitle": "2025寒假高中小组课"
    }


def fill_picture_placeholder(slide, image_source):
    """
    填充图片占位符
    
    参数:
        slide: 幻灯片对象
        image_source: 图片来源，可以是BytesIO对象或文件路径字符串
    """
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if "PICTURE" in str(ph_type):
                # 找到图片占位符，插入图片
                if image_source:
                    try:
                        # 获取占位符位置和大小
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # 删除占位符
                        sp = shape.element
                        sp.getparent().remove(sp)
                        
                        # 在相同位置插入图片
                        slide.shapes.add_picture(image_source, left, top, width, height)
                        print(f"    ✅ 图片已填充到占位符")
                        return True
                    except Exception as e:
                        print(f"    ⚠️ 填充图片失败: {e}")
                        return False
    return False


def get_mindmap_image(data, target_type="learning_objectives"):
    """
    获取思维导图图片路径
    思维导图已在parser.py中精准提取
    
    参数:
        data: 课程数据
        target_type: 目标类型 ("learning_objectives" 或 "summary")
    
    返回:
        图片路径或None
    """
    extracted_images = data.get("extracted_images", [])
    
    if not extracted_images:
        return None
    
    # 查找标记为思维导图的图片
    for img_info in extracted_images:
        if img_info.get("is_mindmap", False):
            img_path = img_info["path"]
            if os.path.exists(img_path):
                print(f"    📊 使用提取的思维导图: {img_info['filename']}")
                return img_path
    
    return None


def generate_ppt():
    """生成PPT主流程"""
    print("=" * 80)
    print("🚀 启动新版PPT生成器（统一模板）")
    print("=" * 80)
    
    # 1. 加载数据
    print("\n[1/4] 加载课程数据...")
    data = load_course_data()
    if not data:
        return
    
    print(f"  ✅ 数据加载成功")
    print(f"  - 讲义标题: {data.get('lecture_title', '未提取')}")
    print(f"  - 学习目标: {len(data.get('learning_objectives', []))} 个")
    print(f"  - 知识点: {len(data.get('knowledge_points', []))} 个")
    
    # 2. 加载模板
    print("\n[2/4] 加载PPT模板...")
    if not os.path.exists(config.MASTER_TEMPLATE):
        print(f"❌ 错误: 找不到模板文件 {config.MASTER_TEMPLATE}")
        return
    
    prs = Presentation(config.MASTER_TEMPLATE)
    print(f"  ✅ 模板加载成功")
    print(f"  - 可用布局: {len(prs.slide_layouts)} 个")
    
    # 删除模板中的预设幻灯片
    if len(prs.slides) > 0:
        print(f"  🗑️ 删除模板中的 {len(prs.slides)} 张预设幻灯片...")
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    
    # 3. 创建幻灯片
    print("\n[3/4] 生成幻灯片...")
    builder = SlideBuilder(prs)
    cover_info = get_cover_info()
    
    slide_count = 0
    
    # ========== 1. 封面（布局0：Cover_Layout）==========
    print("\n  📖 [1] 封面")
    subject = cover_info.get("subject", "语文")
    season = cover_info.get("season", "寒假")
    
    # 生成季节背景图
    cover_bg = generate_cover_image(subject, season)
    
    slide = builder.create_slide(0)
    
    # 添加背景图
    if cover_bg:
        builder.add_background_image(slide, cover_bg)
    
    # 填充4个占位符（保留格式）
    placeholders = list(slide.placeholders)
    for ph in placeholders:
        idx = ph.placeholder_format.idx
        if idx == 10:
            ph.text = f"小组课 · {season}课堂"
        elif idx == 11:
            ph.text = subject
        elif idx == 12:
            ph.text = cover_info.get('subtitle', '2025寒假高中小组课')
        elif idx == 13:
            details = f"高中{subject}·{cover_info.get('grade', '高一')}\n主讲人：{cover_info.get('teacher', 'XXX老师')}"
            ph.text = details
    
    slide_count += 1
    
    # ========== 2. 课程体系（布局1）==========
    print("  📚 [2] 课程体系")
    slide = builder.create_slide(1)
    # 添加课程体系图片（如果有）
    course_system_img = "Smart_PPT_Factory/assets/课程体系.png"
    if os.path.exists(course_system_img):
        builder.add_image(slide, course_system_img, left=2, top=2, width=12, height=6)
    slide_count += 1
    
    # ========== 3. 讲义标题（布局2）- 有图片占位符 ==========
    print("  📝 [3] 讲义标题")
    lecture_title = data.get("lecture_title", "本节课主题")
    slide = builder.create_slide(2)
    
    # 填充标题占位符
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:  # TITLE
            ph.text = lecture_title
    
    # 生成并填充图片占位符
    title_img = generate_lecture_title_image(lecture_title)
    fill_picture_placeholder(slide, title_img)
    slide_count += 1
    
    # ========== 4. 学习目标（布局3）- 标题+图片占位符 ==========
    print("  🎯 [4] 学习目标")
    objectives = data.get("learning_objectives", ["暂无学习目标"])
    slide = builder.create_slide(3)
    
    # 填充标题占位符
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:  # TITLE
            ph.text = "本节课学习目标"
    
    # 使用AI生成学习目标层级图（AI自由创作）
    from ai_image_generator import generate_learning_objectives_image
    objectives_img = generate_learning_objectives_image(objectives)
    
    if objectives_img:
        fill_picture_placeholder(slide, objectives_img)
        print("    ✅ 已生成学习目标层级图（AI创意风格）")
    else:
        print("    ⚠️ 学习目标图生成失败")
    
    slide_count += 1
    
    # ========== 5. 学习目标思维导图（布局4）- 图片占位符 ==========
    mindmap_img = get_mindmap_image(data, "learning_objectives")
    if mindmap_img:
        print("  🗺️ [5] 学习目标思维导图")
        slide = builder.create_slide(4)
        fill_picture_placeholder(slide, mindmap_img)
        slide_count += 1
    
    # ========== 6. 考情（布局5）==========
    print("  📊 [5] 考情分析")
    exam_analysis = data.get("exam_analysis", "暂无考情分析")
    slide = builder.create_slide(5)
    
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "本节课考情"
        elif idx == 11:
            ph.text = exam_analysis
    slide_count += 1
    
    # ========== 知识点循环 ==========
    knowledge_points = data.get("knowledge_points", [])
    
    if not knowledge_points:
        print("\n  ⚠️ 警告: 未找到知识点")
        knowledge_points = [{
            "title": "示例知识点",
            "content": "这是示例内容"
        }]
    
    print(f"\n  📚 知识点部分 ({len(knowledge_points)} 个知识点)")
    
    for i, kp in enumerate(knowledge_points, 1):
        kp_title = kp.get("title", f"知识点{i}")
        kp_content = kp.get("content", "暂无内容")
        
        print(f"\n    知识点 {i}: {kp_title}")
        
        # 6. 知识点切片标题（布局6）- 有图片占位符
        print(f"      [{slide_count+1}] 切片标题")
        slide = builder.create_slide(6)
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 1:
                ph.text = kp_title
        # 生成并填充图片
        kp_title_img = generate_knowledge_point_image(kp_title, "")
        fill_picture_placeholder(slide, kp_title_img)
        slide_count += 1
        
        # 7. 知识点（布局7）- 有图片占位符
        print(f"      [{slide_count+1}] 知识点内容")
        slide = builder.create_slide(7)
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 0:
                ph.text = kp_title
            elif idx == 12:
                ph.text = kp_content
        # 生成并填充图片
        kp_img = generate_knowledge_point_image(kp_title, kp_content[:100])
        fill_picture_placeholder(slide, kp_img)
        slide_count += 1
        
        # 8. 开口说（布局8）- 只在第一个知识点后
        if i == 1:
            discussion = kp.get("discussion", "请思考并讨论相关问题")
            print(f"      [{slide_count+1}] 开口说")
            slide = builder.create_slide(8)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 10:
                    ph.text = discussion
            slide_count += 1
        
        # 9. 经典例题母题（布局9）
        example_mother = kp.get("example_mother", "")
        if example_mother:
            print(f"      [{slide_count+1}] 经典例题（母题）")
            slide = builder.create_slide(9)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 10:
                    ph.text = example_mother
            slide_count += 1
        
        # 10. 经典例题变式（布局10）
        example_variant = kp.get("example_variant", "")
        method = kp.get("method", "")
        if example_variant or method:
            print(f"      [{slide_count+1}] 经典例题（变式/方法）")
            slide = builder.create_slide(10)
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 10:
                    ph.text = example_variant
                elif idx == 11:
                    ph.text = method
            slide_count += 1
    
    # ========== 11. 上台讲（布局11）- 所有知识点完成后 ==========
    print(f"\n  🎤 [{slide_count+1}] 上台讲")
    slide = builder.create_slide(11)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text = "请结合所学知识点，上台分享你的理解和心得"
    slide_count += 1
    
    # ========== 12. 课堂总结过渡（布局12）==========
    print(f"  📋 [{slide_count+1}] 课堂总结过渡")
    slide = builder.create_slide(12)
    slide_count += 1
    
    # ========== 13. 课堂总结内容（布局13）- 有图片占位符 ==========
    print(f"  📋 [{slide_count+1}] 课堂总结内容")
    slide = builder.create_slide(13)
    # 优先使用提取的思维导图，否则AI生成
    summary_mindmap = get_mindmap_image(data, "summary")
    if summary_mindmap:
        print(f"    📊 使用提取的思维导图")
        fill_picture_placeholder(slide, summary_mindmap)
    else:
        summary_img = generate_knowledge_point_image("课堂总结", "本节课重点内容回顾")
        fill_picture_placeholder(slide, summary_img)
    slide_count += 1
    
    # ========== 14. 出门测过渡（布局14）==========
    print(f"  ✅ [{slide_count+1}] 出门测过渡")
    slide = builder.create_slide(14)
    slide_count += 1
    
    # ========== 15. 出门测计时（布局15）==========
    print(f"  ⏱️ [{slide_count+1}] 出门测计时")
    quiz_content = data.get("quiz_content") or "请完成讲义上的测试题"
    slide = builder.create_slide(15)
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = quiz_content
    slide_count += 1
    
    # ========== 16. 作业布置（布局16）==========
    print(f"  📝 [{slide_count+1}] 作业布置")
    homework = data.get("homework", "完成对应练习题")
    slide = builder.create_slide(16)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text = homework
    slide_count += 1
    
    # ========== 17. 告别（布局17）==========
    print(f"  👋 [{slide_count+1}] 告别")
    slide = builder.create_slide(17)
    slide_count += 1
    
    # 4. 保存文件
    print(f"\n[4/4] 保存PPT文件...")
    os.makedirs(os.path.dirname(config.OUTPUT_PATH), exist_ok=True)
    prs.save(config.OUTPUT_PATH)
    
    print("\n" + "=" * 80)
    print(f"✅ PPT生成完成！")
    print(f"📄 文件路径: {config.OUTPUT_PATH}")
    print(f"📊 总页数: {slide_count} 页")
    print("=" * 80)


if __name__ == "__main__":
    generate_ppt()
