"""
分析新增的课堂引入版式
"""
from pptx import Presentation

prs = Presentation("Smart_PPT_Factory/assets/master_template.pptx")

print("=" * 80)
print("PPT模板布局分析")
print("=" * 80)

for i, layout in enumerate(prs.slide_layouts):
    print(f"\n布局 {i}: {layout.name}")
    print(f"  占位符数量: {len(layout.placeholders)}")
    
    for ph in layout.placeholders:
        print(f"    - 索引{ph.placeholder_format.idx}: {ph.placeholder_format.type} ({ph.name})")
        if hasattr(ph, 'text_frame') and ph.text_frame:
            print(f"      文本: {ph.text_frame.text[:50] if ph.text_frame.text else '(空)'}")

print("\n" + "=" * 80)
print(f"总布局数: {len(prs.slide_layouts)}")
print("=" * 80)
