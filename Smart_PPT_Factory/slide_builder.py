"""
增强版幻灯片构建器
支持无占位符布局的文本框添加和智能内容填充
"""
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class SlideBuilder:
    """幻灯片构建器类"""
    
    def __init__(self, prs):
        self.prs = prs
        self.slide_width = prs.slide_width
        self.slide_height = prs.slide_height
    
    def get_layout(self, index):
        """安全获取布局"""
        try:
            return self.prs.slide_layouts[index]
        except IndexError:
            print(f"⚠️ 布局索引 {index} 不存在，使用默认布局")
            return self.prs.slide_layouts[1]
    
    def create_slide(self, layout_index):
        """创建幻灯片"""
        layout = self.get_layout(layout_index)
        return self.prs.slides.add_slide(layout)
    
    def fill_placeholders(self, slide, **kwargs):
        """
        智能填充占位符
        kwargs: title, subtitle, content, body 等
        """
        placeholders = list(slide.placeholders)
        
        # 尝试按名称匹配
        for ph in placeholders:
            ph_name = ph.name.lower()
            
            # 标题
            if 'title' in kwargs and ('标题' in ph_name or 'title' in ph_name):
                ph.text = str(kwargs['title'])
            
            # 副标题
            elif 'subtitle' in kwargs and ('副标题' in ph_name or 'subtitle' in ph_name):
                ph.text = str(kwargs['subtitle'])
            
            # 内容/正文
            elif 'content' in kwargs and ('内容' in ph_name or 'content' in ph_name or 'object' in ph_name):
                ph.text = str(kwargs['content'])
            
            # Body
            elif 'body' in kwargs and ('正文' in ph_name or 'body' in ph_name or '文本' in ph_name):
                ph.text = str(kwargs['body'])
        
        # 如果没有匹配到，按索引填充
        if len(placeholders) > 0 and 'title' in kwargs:
            placeholders[0].text = str(kwargs['title'])
        
        if len(placeholders) > 1 and 'content' in kwargs:
            placeholders[1].text = str(kwargs['content'])
    
    def add_textbox(self, slide, text, left, top, width, height, 
                    font_size=28, font_name="微软雅黑", bold=False, 
                    color=(0, 0, 0), align="left"):
        """
        添加文本框
        
        参数:
            slide: 幻灯片对象
            text: 文本内容
            left, top, width, height: 位置和尺寸（英寸）
            font_size: 字号
            font_name: 字体
            bold: 是否加粗
            color: RGB颜色元组
            align: 对齐方式 (left/center/right)
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.text = str(text)
        
        for paragraph in tf.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = RGBColor(*color)
            
            if align == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT
        
        return textbox
    
    def add_title(self, slide, text, font_size=40):
        """添加标题（标准位置）"""
        return self.add_textbox(
            slide, text,
            left=0.5, top=0.5, width=15, height=1.2,
            font_size=font_size, bold=True
        )
    
    def add_content(self, slide, text, font_size=24):
        """添加正文内容（标准位置）"""
        return self.add_textbox(
            slide, text,
            left=1, top=2, width=14, height=6,
            font_size=font_size
        )
    
    def add_centered_title(self, slide, text, font_size=48):
        """添加居中标题"""
        return self.add_textbox(
            slide, text,
            left=2, top=3.5, width=12, height=2,
            font_size=font_size, bold=True, align="center"
        )
    
    def add_background_image(self, slide, image_path_or_stream):
        """添加背景图片并置底"""
        try:
            pic = slide.shapes.add_picture(
                image_path_or_stream, 
                0, 0, 
                width=self.slide_width, 
                height=self.slide_height
            )
            
            # 置底
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            
            return pic
        except Exception as e:
            print(f"⚠️ 添加背景图片失败: {e}")
            return None
    
    def add_image(self, slide, image_path_or_stream, left, top, width=None, height=None):
        """添加图片（支持文件路径或BytesIO流）"""
        try:
            if width and height:
                pic = slide.shapes.add_picture(
                    image_path_or_stream,
                    Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            else:
                pic = slide.shapes.add_picture(
                    image_path_or_stream,
                    Inches(left), Inches(top)
                )
            return pic
        except Exception as e:
            print(f"⚠️ 添加图片失败: {e}")
            return None


def create_cover_slide_with_master(cover_prs, cover_info, bg_image=None):
    """
    使用Cover_Layout母版创建封面
    
    参数:
        cover_prs: 封面模板的Presentation对象
        cover_info: 封面信息字典
        bg_image: 背景图片（可选）
    
    返回:
        创建的封面slide
    """
    # 使用Cover_Layout布局（索引0）
    cover_layout = cover_prs.slide_layouts[0]
    slide = cover_prs.slides.add_slide(cover_layout)
    
    # 添加背景图（如果有）- 放在最底层，不遮挡文字
    if bg_image:
        pic = slide.shapes.add_picture(
            bg_image, 0, 0,
            width=cover_prs.slide_width,
            height=cover_prs.slide_height
        )
        # 置底
        try:
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        except:
            pass
    
    # 填充4个占位符
    # 根据PDF文件名解析的信息填充
    placeholders = list(slide.placeholders)
    
    if len(placeholders) >= 4:
        # 占位符[10]: 小组课·X季课堂
        season = cover_info.get("season", "寒假")
        placeholders[0].text = f"小组课 · {season}课堂"
        
        # 占位符[11]: 主科目（大字，如"语文"）
        subject = cover_info.get("subject", "语文")
        placeholders[1].text = subject
        
        # 占位符[12]: 副标题（如"2025寒假高中小组课"）
        subtitle = cover_info.get("subtitle", "2025寒假高中小组课")
        placeholders[2].text = subtitle
        
        # 占位符[13]: 详细信息（如"高中语文·高一 主讲人：XXX老师"）
        details = f"高中{subject}·{cover_info.get('grade', '高一')} 主讲人：{cover_info.get('teacher', 'XXX老师')}"
        placeholders[3].text = details
    
    return slide


def create_course_system_slide(builder):
    """创建课程体系页（布局19：开场视频）"""
    slide = builder.create_slide(19)
    
    # 添加标题
    builder.add_textbox(slide, "志高班-语文课程体系",
                       left=0.5, top=0.5, width=15, height=1,
                       font_size=36, bold=True)
    
    # 添加课程体系图片
    image_path = "Smart_PPT_Factory/assets/课程体系.png"
    if os.path.exists(image_path):
        builder.add_image(slide, image_path, left=2, top=2, width=12, height=6)
    else:
        print(f"⚠️ 未找到课程体系图片: {image_path}")
    
    return slide


def create_lecture_title_slide(builder, title, title_image=None):
    """创建讲义标题页（布局14：讲义标题）"""
    slide = builder.create_slide(14)
    builder.fill_placeholders(slide, title=title)
    
    # 添加标题配图（如果有）
    if title_image:
        # 放在右侧作为装饰
        builder.add_image(slide, title_image, left=12, top=2, width=3, height=3)
    
    return slide


def create_learning_objectives_slide(builder, objectives):
    """创建学习目标页（布局1：标题和内容）"""
    slide = builder.create_slide(1)
    
    # 格式化学习目标
    objectives_text = "\n".join([f"• {obj}" for obj in objectives])
    
    builder.fill_placeholders(
        slide,
        title="本节课学习目标",
        content=objectives_text
    )
    
    return slide


def create_exam_analysis_slide(builder, analysis):
    """创建考情分析页（布局1：标题和内容）"""
    slide = builder.create_slide(1)
    builder.fill_placeholders(
        slide,
        title="本节课考情",
        content=analysis
    )
    return slide


def create_knowledge_section_title_slide(builder, title):
    """创建知识点切片标题页（布局20：知识点切片标题）"""
    slide = builder.create_slide(20)
    builder.fill_placeholders(slide, title=title)
    return slide


def create_knowledge_checklist_slide(builder, title, content):
    """创建知识清单页（布局24：知识清单+笔记页）"""
    slide = builder.create_slide(24)
    builder.fill_placeholders(
        slide,
        title=title,
        content=content
    )
    return slide


def create_knowledge_lecture_slide(builder, title, content):
    """创建知识点讲解页（布局18：知识点互动+讲解页）"""
    slide = builder.create_slide(18)
    builder.fill_placeholders(slide, title=title)
    
    # 布局18只有标题占位符，内容需要手动添加
    builder.add_content(slide, content, font_size=24)
    
    return slide


def create_discussion_slide(builder, question):
    """创建开口说页（布局11：开口说）"""
    slide = builder.create_slide(11)
    builder.fill_placeholders(slide, content=question)
    return slide


def create_example_slide(builder, example_content):
    """创建经典例题页（布局25：经典例题）"""
    slide = builder.create_slide(25)
    
    # 布局25没有占位符，需要手动添加
    builder.add_textbox(slide, "经典例题",
                       left=0.5, top=0.5, width=6, height=0.8,
                       font_size=32, bold=True, color=(0, 102, 204))
    
    builder.add_textbox(slide, example_content,
                       left=1, top=1.5, width=14, height=6.5,
                       font_size=22)
    
    return slide


def create_method_summary_slide(builder, method_content):
    """创建解题方法总结页（布局17：解题方法总结）"""
    slide = builder.create_slide(17)
    builder.fill_placeholders(slide, body=method_content)
    return slide


def create_section_title_slide(builder, title):
    """创建章节标题页（布局15：2_标题幻灯片）"""
    slide = builder.create_slide(15)
    
    # 布局15没有占位符，需要居中添加标题
    builder.add_centered_title(slide, title, font_size=56)
    
    return slide


def create_summary_slide(builder):
    """创建总结页（布局21：总结页）"""
    slide = builder.create_slide(21)
    # 布局21通常是图片页，可能需要添加总结思维导图
    return slide


def create_homework_slide(builder, homework):
    """创建作业页（布局10：课后作业）"""
    slide = builder.create_slide(10)
    builder.fill_placeholders(slide, body=homework)
    return slide


def create_ending_slide(builder):
    """创建结束页（布局16：1_结束页）"""
    slide = builder.create_slide(16)
    # 布局16通常是固定设计的结束页
    return slide
